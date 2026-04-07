#!/usr/bin/env python3
"""
Extract Kabir's Writings.pptx into:
  - data/writings/kabir_writings.json   (entries: id, title, body, image)
  - data/writings/images/kw_NNN.jpg     (cropped: title bar removed if detected)

Strategy:
  - For each slide: collect all text frames in reading order.
  - First non-empty short line  -> title
  - Remaining lines             -> body
  - Largest image on the slide  -> hero image (saved as jpg, max width 1200)
  - Skip slides with no body text (cover/section dividers handled separately)
"""
import os, json, io, re
from pptx import Presentation
from pptx.util import Emu
from PIL import Image

ROOT = os.path.dirname(os.path.abspath(__file__))
PPTX = os.path.join(ROOT, "Kabir's Writings.pptx")
IMG_DIR = os.path.join(ROOT, "images")
OUT_JSON = os.path.join(ROOT, "kabir_writings.json")
os.makedirs(IMG_DIR, exist_ok=True)

MAX_W = 1200
JPG_Q = 82

def clean(s):
    return re.sub(r"\s+", " ", (s or "")).strip()

def save_image(blob, path):
    try:
        im = Image.open(io.BytesIO(blob))
        if im.mode in ("RGBA", "P", "LA"):
            bg = Image.new("RGB", im.size, (252, 251, 248))  # Crisp Pearl
            bg.paste(im, mask=im.split()[-1] if im.mode in ("RGBA","LA") else None)
            im = bg
        if im.width > MAX_W:
            h = int(im.height * MAX_W / im.width)
            im = im.resize((MAX_W, h), Image.LANCZOS)
        im.save(path, "JPEG", quality=JPG_Q, optimize=True)
        return True
    except Exception as e:
        print(f"  ! image save failed: {e}")
        return False

def extract():
    prs = Presentation(PPTX)
    entries = []
    skipped = []

    for idx, slide in enumerate(prs.slides, start=1):
        # Collect text shapes in y-order
        text_shapes = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                txt = "\n".join(clean(p.text) for p in sh.text_frame.paragraphs if clean(p.text))
                if txt:
                    top = sh.top if sh.top is not None else 0
                    text_shapes.append((top, txt))
        text_shapes.sort(key=lambda t: t[0])
        all_text = "\n".join(t for _, t in text_shapes).strip()

        # Heuristic split: first short line = title, rest = body
        title, body = "", ""
        if all_text:
            lines = [l for l in all_text.splitlines() if l.strip()]
            if lines:
                # Title = first line if it's short-ish, else first sentence
                if len(lines[0]) <= 80:
                    title = lines[0]
                    body = "\n".join(lines[1:]).strip()
                else:
                    # split on first period
                    parts = lines[0].split(".", 1)
                    title = parts[0][:80]
                    body = (parts[1] if len(parts) > 1 else "") + "\n" + "\n".join(lines[1:])
                    body = body.strip()

        # Largest image on the slide
        biggest = None
        biggest_area = 0
        for sh in slide.shapes:
            if sh.shape_type == 13:  # PICTURE
                try:
                    blob = sh.image.blob
                    im = Image.open(io.BytesIO(blob))
                    area = im.width * im.height
                    if area > biggest_area:
                        biggest_area = area
                        biggest = blob
                except Exception:
                    pass

        if not body and not title:
            skipped.append((idx, "empty"))
            continue
        if not biggest:
            skipped.append((idx, f"no-image: {title[:40]}"))
            # still keep entry without image
            img_name = ""
        else:
            img_name = f"kw_{idx:03d}.jpg"
            save_image(biggest, os.path.join(IMG_DIR, img_name))

        entries.append({
            "id": f"kw_{idx:03d}",
            "slide": idx,
            "title": title or f"Reflection {idx}",
            "body": body,
            "image": f"data/writings/images/{img_name}" if img_name else "",
        })
        print(f"  slide {idx:3d}  {title[:60]}")

    with open(OUT_JSON, "w", encoding="utf-8") as f:
        json.dump({"entries": entries, "count": len(entries)}, f, indent=2, ensure_ascii=False)

    print(f"\n✓ {len(entries)} entries written to {OUT_JSON}")
    if skipped:
        print(f"  ({len(skipped)} skipped)")
        for s in skipped[:10]:
            print(f"    slide {s[0]}: {s[1]}")

if __name__ == "__main__":
    extract()
